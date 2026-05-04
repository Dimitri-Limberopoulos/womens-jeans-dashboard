#!/usr/bin/env python3
"""
build_deck.py — Bharat-style deck for the Target Owned Brand team.

8 collapsed insights (down from 13) + a grouped summary cover.
Each insight slide uses 1-3 full-width chart placeholders so charts can
be dropped in directly via ThinkCell without cramped real estate.

Layout per insight slide:
  - Title bar (Insight NN | Title)
  - One-line bold takeaway
  - Section divider
  - Chart area (1, 2, or 3 placeholders, each with chart title + brief)
  - Single-line italic data summary above the footer
  - Footer + page number
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── COLORS ────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY        = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY   = RGBColor(0x55, 0x55, 0x55)
SUBTITLE    = RGBColor(0x44, 0x44, 0x44)
LIGHT_GRAY  = RGBColor(0xC0, 0xC0, 0xC0)
CARD_BG     = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT      = RGBColor(0xCC, 0x00, 0x00)  # Target red
PLACEHOLDER_BG = RGBColor(0xFA, 0xFA, 0xFA)

FONT_HDR    = "Arial Black"
FONT_BODY   = "Arial"

# ─── DIMENSIONS ────────────────────────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
CONTENT_W = Inches(9.0)


# ─── HELPERS ───────────────────────────────────────────────────────────────

def set_slide_size(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_text_run(paragraph, text, *, font=FONT_BODY, size=11, bold=False,
                 italic=False, color=BLACK):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, anchor="top"):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    return tf


def add_title_bar(slide, bold_part, light_part, y=Inches(0.25)):
    tf = add_textbox(slide, MARGIN_L, y, CONTENT_W, Inches(0.4))
    p = tf.paragraphs[0]
    add_text_run(p, bold_part, font=FONT_HDR, size=18, bold=True, color=BLACK)
    add_text_run(p, "  |  ", font=FONT_BODY, size=18, color=LIGHT_GRAY)
    add_text_run(p, light_part, font=FONT_BODY, size=18, color=SUBTITLE)


def add_divider(slide, y, color=LIGHT_GRAY, weight=0.5):
    line = slide.shapes.add_connector(
        1, MARGIN_L, y, MARGIN_L + CONTENT_W, y,
    )
    line.line.color.rgb = color
    line.line.width = Pt(weight)


def _kill_shadow(shape):
    sppr = shape.fill._xPr
    for el in sppr.findall(qn("a:effectLst")):
        sppr.remove(el)
    etree.SubElement(sppr, qn("a:effectLst"))


def add_chart_placeholder(slide, x, y, w, h, label, sub):
    """Dashed-border placeholder where a ThinkCell chart will go."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER_BG
    box.line.color.rgb = LIGHT_GRAY
    box.line.width = Pt(0.75)
    box.line.dash_style = 7  # dashed
    _kill_shadow(box)

    # Label band — small text up top
    tf = add_textbox(slide, x + Inches(0.1), y + Inches(0.08),
                     w - Inches(0.2), Inches(0.22))
    p = tf.paragraphs[0]
    add_text_run(p, "CHART PLACEHOLDER", font=FONT_HDR, size=8, bold=True,
                 color=ACCENT)

    # Title + brief in middle
    tf2 = add_textbox(slide, x + Inches(0.2), y + Inches(0.35),
                      w - Inches(0.4), h - Inches(0.55), anchor="middle")
    p1 = tf2.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_text_run(p1, label, font=FONT_BODY, size=12, bold=True, color=BLACK)
    if sub:
        p2 = tf2.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        add_text_run(p2, sub, font=FONT_BODY, size=9.5, italic=True, color=GRAY)


def add_page_number(slide, num):
    tf = add_textbox(slide, Inches(9.0), Inches(5.32), Inches(0.9), Inches(0.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_text_run(p, f"{num}", font=FONT_BODY, size=8, color=LIGHT_GRAY)


def add_footer_label(slide, text):
    tf = add_textbox(slide, MARGIN_L, Inches(5.32), Inches(7), Inches(0.2))
    p = tf.paragraphs[0]
    add_text_run(p, text, font=FONT_BODY, size=8, color=LIGHT_GRAY)


# ─── INSIGHT CONTENT (8 collapsed insights, renumbered for narrative flow) ─

INSIGHTS = [
    # ── A. ASSORTMENT (Scale + Mix combined) ──────────────────────────────
    {
        "num": 1,
        "group": "Assortment",
        "title": "OB Assortment Lags All Competitors",
        "summary_desc": (
            "95 CCs across 3 brands — the smallest owned-brand jeans "
            "assortment studied. Walmart OB has 4.3× more (406), Kohl's 520, "
            "Macy's 335. Depth/style is also thin: 2.2 colors/style vs "
            "Walmart 3.6 and Old Navy 4.2."
        ),
        "headline": "Target OB is the smallest owned-brand jeans assortment of any retailer studied — and the plus-size sub-brand is materially undersized.",
        "summary": "Total CCs: Target 95 · Walmart 406 · Macy's 335 · Kohl's 520 · Amazon 78. Plus-OB CCs: Ava & Viv 14 · Terra & Sky 93 · Kohl's SO 166. Depth/style: Target 2.2 vs Walmart 3.6, Old Navy 4.2, Amazon 4.6.",
        "charts": [
            ("Total OB Assortment by Retailer (CCs)",
             "Horizontal bar, sorted descending. Target OB highlighted. Annotate # styles + colors/style."),
            ("Plus-Focused OB Depth — CC Count",
             "Bar comparison: Ava & Viv vs Terra & Sky vs Kohl's SO. Highlight the 14 vs 93 vs 166 gap."),
        ],
    },
    {
        "num": 2,
        "group": "Assortment",
        "title": "Silhouette + Rise Mix Misaligned",
        "summary_desc": (
            "Wide Leg 12% vs ~26% competitor avg — under-indexed on the "
            "dominant trend silhouette. Low Rise 23% (Wild Fable–driven) "
            "vs <11% at Walmart, Macy's, Kohl's — over-exposed to a "
            "fading rise category. Trend-cycle risk concentrated."
        ),
        "headline": "Target OB under-indexes on the dominant Wide Leg silhouette and over-indexes on Low Rise — concentrating trend-cycle risk.",
        "summary": "Wide Leg: Target 12% vs Walmart 26%, Macy's 27%, Old Navy 21%, AE 18%. Low Rise: Target 23% vs Walmart 1.5%, Macy's 0%, Kohl's 11%, AE 26%. Wild Fable's low-rise-heavy assortment becomes hard to move if the trend cycle shifts.",
        "charts": [
            ("Wide Leg Share by OB Group",
             "Bar chart sorted descending. Target OB highlighted as laggard at 12%."),
            ("Low Rise Share by OB Group",
             "Bar chart sorted descending. Target OB elevated at 23%; only AE (26%) higher. Annotate Wild Fable as driver."),
        ],
    },
    {
        "num": 3,
        "group": "Assortment",
        "title": "Wash Mix Skews Light",
        "summary_desc": (
            "Light Wash 23% — highest of any OB. Dark Wash 16% — below "
            "Walmart 20% and Kohl's 28%. Inverted vs the cross-shop "
            "pattern. Rebalancing toward dark and medium could lift "
            "year-round sell-through and cross-shopping appeal."
        ),
        "headline": "Target OB indexes highest in Light Wash and lightest in Dark — opposite the broader OB pattern.",
        "summary": "Light Wash: Target 23% · Walmart 18% · Kohl's 18% · Macy's 8%. Dark Wash: Target 16% · Walmart 20% · Macy's 17% · Kohl's 28%. Medium Wash broadly aligned at ~27%.",
        "charts": [
            ("Wash Mix by OB Group (100% Stacked)",
             "Horizontal 100% stacked bar — Light / Medium / Dark / Black / Color. One row per OB group. Light slice prominent for Target; Dark slice prominent for Kohl's. Full width for clarity."),
        ],
    },
    # ── B. PRICING & PROMO ────────────────────────────────────────────────
    {
        "num": 4,
        "group": "Pricing & Promo",
        "title": "Three Brands, One Tight Price Cluster",
        "summary_desc": (
            "58% of OB CCs in a single $28–$35 band. Universal Thread "
            "$28, Wild Fable $32, Ava & Viv $33 stack on top of each "
            "other. No premium tier above $40, no entry tier below $20 — "
            "three brands compete for the same basket."
        ),
        "headline": "Three OB brands stack on top of each other in a single $28-$35 band — no upmarket or downmarket ladder.",
        "summary": "Universal Thread median $28 (44 CCs) · Wild Fable $32 (37) · Ava & Viv $33 (14). Combined IQR $28-$36. No premium tier above $40, no entry tier below $20.",
        "charts": [
            ("Target OB Price Distribution by Brand",
             "Box plot per brand at current price. Overlay shaded $28-$35 cluster band. Single chart, full-width."),
            ("Share of Target OB CCs by Price Band",
             "Horizontal stacked bar: <$20 / $20-$28 / $28-$35 / $35-$45 / $45+. The 58% in $28-$35 should dominate."),
        ],
    },
    {
        "num": 5,
        "group": "Pricing & Promo",
        "title": "OB Price Pressed from All Sides at Current",
        "summary_desc": (
            "Macy's OB collapses $60 → $36 (90% on sale, 40% off). "
            "Kohl's compresses to $30–$40 (71% on sale). Target's own 3P "
            "at $35–$50 IQR (87% on sale) — 30% sits at or below OB Q3 "
            "($35). Only Target NB (1P) holds firm at $65 med."
        ),
        "headline": "Competitor discounting and Target's own 3P erode the OB premium read; only Target NB (1P) holds price firmly.",
        "summary": "Macy's OB collapses $60 med → $36 (90% on sale, 40% off). Kohl's compresses to $30-$40 (71% on sale). Target 3P at $35-$50 (87% on sale, 34% off) — 30% below OB Q3 ($35). Target NB stays at $65 med, only 13% on sale.",
        "charts": [
            ("Original vs Current Price IQR — OB Groups",
             "Floating-bar comparison per OB. Show the IQR collapse from original to current. Highlight Macy's $60→$36 collapse."),
            ("Target OB vs NB vs 3P — Current Price Distribution",
             "Three-way box plot (or violin). Annotate the 30% of 3P that sits at or below OB Q3 ($35)."),
        ],
    },
    {
        "num": 6,
        "group": "Pricing & Promo",
        "title": "Promo: Infrequent But Deep — an Outlier",
        "summary_desc": (
            "20% on sale at avg 43% off. Walmart 14% / Old Navy 4% (low "
            "frequency). Kohl's 71% on sale at just 18% off (high "
            "frequency, shallow). Target alone in the deep-infrequent "
            "quadrant — risks training markdown-wait behavior."
        ),
        "headline": "Target OB sits alone on the depth-of-discount axis — risks training the customer to wait for markdowns.",
        "summary": "Target OB 20% on sale at 43% off (deep, infrequent). Walmart OB 14% on sale (low frequency). Old Navy 4% (almost no promo). Kohl's OB 71% on sale at 18% off (high frequency, shallow). Target's posture is the outlier.",
        "charts": [
            ("Promo Posture Across OB Programs",
             "Scatter: x = % CCs on sale, y = avg discount when on sale. Each retailer plotted; Target OB labelled and highlighted as the outlier in the deep-infrequent quadrant."),
        ],
    },
    # ── C. WHITE SPACE ────────────────────────────────────────────────────
    {
        "num": 7,
        "group": "White Space",
        "title": "Plus-Size Customer Served by 3P, Not OB",
        "summary_desc": (
            "50% of Target 3P (1,261 CCs) is plus-size brands — Woman "
            "Within 729, Roaman's 243, Jessica London 106. Ava & Viv has "
            "14 CCs to compete. Walmart's plus posture is materially "
            "deeper: Terra & Sky 93 OB plus Just My Size 214 NB. Largest "
            "OB white space by far."
        ),
        "headline": "Half of Target 3P is plus-size brands. Ava & Viv competes with 14 CCs against 1,261 — the single largest white space for the OB team.",
        "summary": "Target 3P plus brands: Woman Within 729 · Roaman's 243 · Jessica London 106 · AVENUE 67 · Catherines 66. Walmart's plus posture: Terra & Sky 93 OB + Just My Size 214 NB + others.",
        "charts": [
            ("Target Plus-Size Assortment by Source",
             "Stacked bar showing OB / NB / 3P share of plus-size CCs at Target. The 14 vs 1,261 contrast carries the slide."),
            ("Plus-Size Brand Depth Across Retailers",
             "Bar chart of CC counts for the largest plus-focused brand at each retailer (Ava & Viv vs Terra & Sky vs Kohl's SO vs Walmart Just My Size, etc.)."),
        ],
    },
    {
        "num": 8,
        "group": "White Space",
        "title": "NB Strategy: Curated Premium vs Broad Value",
        "summary_desc": (
            "Target NB: 67 CCs / 2 brands (Levi's + KBB). Walmart NB: "
            "2,666 / 36 (Lee 732, Sofia Vergara 274, Gloria Vanderbilt, "
            "Wrangler, Levi Strauss Signature). Levi's plays three "
            "roles: $70 Target (premium) · $40 Walmart (value) · $75 "
            "standalone."
        ),
        "headline": "Target keeps the 1P NB shelf premium and narrow; Walmart competes on breadth. Levi's plays three different roles by retailer.",
        "summary": "Target NB: 67 CCs · 2 brands (Levi's + KBB by Kahlana). Walmart NB: 2,666 CCs · 36 brands (Lee, Sofia Vergara, Gloria Vanderbilt, Wrangler, Jessica Simpson, Levi's, Levi Strauss Signature). Levi's: Target med $70 (15% on sale) · Walmart med $40 (14%) · Standalone med $75 (54%).",
        "charts": [
            ("National-Brand Breadth by Retailer",
             "Two side-by-side bars per retailer: total NB CCs (left) and # of distinct brands (right). Target's narrowness vs Walmart's breadth is the takeaway."),
            ("Levi's Price Posture by Retailer",
             "Three-bar comparison of Levi's median price + sale% across Target NB, Walmart NB, Levi's standalone. Reveals Target's $70 ceiling and Walmart's $40 floor."),
        ],
    },
]


# ─── COVER / SUMMARY ───────────────────────────────────────────────────────

# Three-column cover: Assortment | Pricing & Promo | White Space.
# Each item is either an int (= insight #, pulled from INSIGHTS) or a dict
# overriding the rendered title + description for a combined entry.
SUMMARY_GROUPS = [
    ("A", "Assortment",      [1, 2, 3]),
    ("B", "Pricing & Promo", [4, 5, 6]),
    ("C", "White Space",     [7, 8]),
]


def build_summary(slide):
    add_title_bar(slide,
                  "Cross-Retailer Jeans Analysis",
                  "Key Findings — Target OB Team")
    add_divider(slide, Inches(0.72))

    # Sub-headline
    tf = add_textbox(slide, MARGIN_L, Inches(0.78), CONTENT_W, Inches(0.28))
    p = tf.paragraphs[0]
    add_text_run(p, "8 findings across three themes  ·  ",
                 font=FONT_BODY, size=10, italic=True, color=GRAY)
    add_text_run(p, "Assortment  ·  Pricing & Promo  ·  White Space",
                 font=FONT_BODY, size=10, italic=True, bold=True, color=DARK_GRAY)

    by_num = {ins["num"]: ins for ins in INSIGHTS}

    # Three columns evenly distributed across content width.
    col_gap = Inches(0.25)
    col_w = Inches((9.0 - 0.5) / 3)   # ~2.83"
    col_xs = [
        MARGIN_L,
        MARGIN_L + col_w + col_gap,
        MARGIN_L + 2 * (col_w + col_gap),
    ]
    section_header_h = Inches(0.30)
    y_start = Inches(1.18)
    title_h = Inches(0.30)
    desc_h = Inches(0.92)
    row_total = title_h + desc_h + Inches(0.08)  # = 1.30"

    def resolve_entry(entry):
        """Normalize an entry into (badge_label, title, summary_desc)."""
        if isinstance(entry, int):
            ins = by_num[entry]
            return (f"{entry:02d}", ins["title"], ins["summary_desc"])
        # Combined / overridden entry
        label = entry.get("label") or "·".join(f"{n:02d}" for n in entry.get("nums", []))
        return (label, entry["title"], entry["summary_desc"])

    def render_column(x, letter, header, items):
        tf = add_textbox(slide, x, y_start, col_w, section_header_h)
        p = tf.paragraphs[0]
        add_text_run(p, f"{letter}.  ", font=FONT_HDR, size=11, bold=True, color=ACCENT)
        add_text_run(p, header.upper(), font=FONT_HDR, size=11, bold=True, color=BLACK)
        line_y = y_start + section_header_h - Inches(0.04)
        line = slide.shapes.add_connector(1, x, line_y, x + col_w, line_y)
        line.line.color.rgb = LIGHT_GRAY
        line.line.width = Pt(0.5)

        ry = y_start + section_header_h + Inches(0.10)
        for entry in items:
            label, title, desc = resolve_entry(entry)
            tf1 = add_textbox(slide, x, ry, col_w, title_h)
            tf1.word_wrap = True
            p1 = tf1.paragraphs[0]
            add_text_run(p1, f"{label}  ", font=FONT_HDR, size=10, bold=True, color=ACCENT)
            add_text_run(p1, title, font=FONT_BODY, size=9.5,
                         bold=True, color=BLACK)
            tf2 = add_textbox(slide, x, ry + title_h, col_w, desc_h)
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            add_text_run(p2, desc, font=FONT_BODY, size=8.5, color=GRAY)
            ry += row_total

    for i, (letter, header, items) in enumerate(SUMMARY_GROUPS):
        render_column(col_xs[i], letter, header, items)

    # Footer note
    tf = add_textbox(slide, MARGIN_L, Inches(5.05), CONTENT_W, Inches(0.22))
    p = tf.paragraphs[0]
    add_text_run(p, "Numbers reflect 7,773 colorway records across 11 "
                 "retailer groups (Apr 2026).",
                 font=FONT_BODY, size=8.5, italic=True, color=GRAY)
    add_footer_label(slide, "A&M | Target OB Competitive Analysis")
    add_page_number(slide, 1)


# ─── INSIGHT SLIDE BUILDER ─────────────────────────────────────────────────

def build_insight(slide, ins, page):
    add_title_bar(slide, f"Insight {ins['num']:02d}", ins["title"])
    add_divider(slide, Inches(0.72))

    # Headline takeaway (bold, full width)
    tf = add_textbox(slide, MARGIN_L, Inches(0.85), CONTENT_W, Inches(0.55))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, ins["headline"], font=FONT_BODY, size=11.5, bold=True, color=BLACK)

    # Section divider below takeaway
    add_divider(slide, Inches(1.50))

    # Chart area
    chart_y = Inches(1.65)
    chart_h = Inches(3.20)
    charts = ins["charts"]
    n = len(charts)
    if n == 1:
        # Full-width single chart
        add_chart_placeholder(slide, MARGIN_L, chart_y, CONTENT_W, chart_h,
                              charts[0][0], charts[0][1])
    elif n == 2:
        # Two side-by-side, each ~half width
        gap = Inches(0.20)
        each_w = Inches((9.0 - 0.20) / 2)
        for i, (title, brief) in enumerate(charts):
            cx = MARGIN_L + i * (each_w + gap)
            add_chart_placeholder(slide, cx, chart_y, each_w, chart_h, title, brief)
    elif n == 3:
        # Three side-by-side
        gap = Inches(0.18)
        each_w = Inches((9.0 - 2 * 0.18) / 3)
        for i, (title, brief) in enumerate(charts):
            cx = MARGIN_L + i * (each_w + gap)
            add_chart_placeholder(slide, cx, chart_y, each_w, chart_h, title, brief)

    # Italic data summary just above the footer
    tf = add_textbox(slide, MARGIN_L, Inches(4.95), CONTENT_W, Inches(0.30))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, ins["summary"], font=FONT_BODY, size=8.5,
                 italic=True, color=DARK_GRAY)

    add_footer_label(slide, "A&M | Target OB Competitive Analysis")
    add_page_number(slide, page)


# ─── MAIN ──────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    set_slide_size(prs)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    build_summary(s)

    for i, ins in enumerate(INSIGHTS):
        s = prs.slides.add_slide(blank)
        build_insight(s, ins, page=i + 2)

    out = "Target OB - Cross-Retailer Jeans Insights.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}  ({len(INSIGHTS)} insights + cover)")
    for ins in INSIGHTS:
        print(f"  {ins['num']:02d}  {ins['title']}  ({len(ins['charts'])} chart{'s' if len(ins['charts'])>1 else ''})")


if __name__ == "__main__":
    main()
